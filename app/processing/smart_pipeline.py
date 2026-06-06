"""
Smart Pipeline Orchestrator

Processes PDF and PPTX files into clean Markdown using local font-aware
extraction, heuristic merging, and an optional AI polish pass.
"""

import os
import re
import logging
from pathlib import Path
from typing import Dict, List, Optional, Callable, Any

from app.processing.font_extractor import FontAwareExtractor
from app.processing.signal_merger import SignalMerger, blocks_to_markdown
from app.config import get_settings

logger = logging.getLogger(__name__)
settings = get_settings()


class SmartPipeline:
    """
    Orchestrates the multi-method extraction pipeline.

    Args:
        use_layout_detection: Legacy flag, kept for backward-compat (unused)
        use_table_transformer: Legacy flag, kept for backward-compat (unused)
        gemini_api_key: Gemini API key for AI polish; if None, polish is disabled
        gemini_model: Model to use for the AI polish pass
        use_polish: When True, use Gemini to polish the final markdown
    """

    def __init__(
        self,
        use_layout_detection: bool = False,
        use_table_transformer: bool = False,
        gemini_api_key: Optional[str] = None,
        gemini_model: Optional[str] = None,
        use_polish: bool = False,
        # Legacy kwargs accepted but ignored
        use_vision: bool = False,
        inter_call_delay_s: float = 0.0,
    ):
        self.use_polish = use_polish and (bool(gemini_api_key) or bool(settings.GLOBAL_AI_TIER1_API_KEY))
        self.gemini_api_key = gemini_api_key
        self.gemini_model = gemini_model or settings.GLOBAL_AI_TIER1_MODEL
        self.font_extractor = FontAwareExtractor()
        self.layout_detector = None  # Legacy: disabled
        self.table_detector = None   # Legacy: disabled
        self.merger = SignalMerger()

    def process(self, file_path: str, progress_callback: Optional[Callable[[int], None]] = None) -> str:
        """
        Process a PDF or PPTX file and return clean Markdown.
        """
        file_path = str(file_path)
        ext = Path(file_path).suffix.lower()

        # Check for empty/missing file
        if not os.path.exists(file_path) or os.path.getsize(file_path) == 0:
            return f"Error: File is missing or empty: {file_path}"

        markdown = ""
        try:
            if progress_callback:
                progress_callback(10) # 10%: Started local process

            markdown = self._local_process(file_path, ext)

            # Final AI Polish Pass
            if self.use_polish and markdown:
                if progress_callback:
                    progress_callback(30) # 30%: Local done, starting AI
                markdown = self._ai_polish(markdown, progress_callback=progress_callback)

            if progress_callback:
                progress_callback(100) # 100%: All done

            # Quality metrics (final)
            lines = [l for l in markdown.split("\n") if l.strip()]
            headings = len([l for l in lines if l.startswith("#")])
            list_items = len([l for l in lines if l.startswith("- ") or l.startswith("1. ")])
            logger.info(f"Final Output: {len(lines)} lines, {headings} headings, {list_items} list items")

            return markdown
        except Exception as e:
            logger.error(f"Failed to process {file_path}: {e}", exc_info=True)
            return f"Error: Could not process file due to extraction failure: {e}"

    def _local_process(self, file_path: str, ext: str) -> str:
        """Internal helper for local extraction."""
        if ext == ".pdf":
            return self._process_pdf(file_path)
        elif ext == ".pptx":
            return self._process_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")


    def _process_pdf(self, pdf_path: str) -> str:
        """Process a PDF file through the pipeline."""
        logger.info(f"Processing PDF: {pdf_path}")
        import pdfplumber

        # Extract tables first (for position tracking)
        logger.info("Extracting tables...")
        tables = self._extract_tables_from_pdf(pdf_path)
        logger.info(f"  Found tables on {len([t for t in tables if t])} pages")

        # Font-aware extraction (primary method)
        logger.info("Extracting text via font-aware method...")
        font_results = self.font_extractor.extract(pdf_path, table_bboxes_per_page={})
        logger.info(f"  Extracted {sum(len(p['blocks']) for p in font_results)} blocks from {len(font_results)} pages")

        # Merge signals
        logger.info("Merging signals...")
        merged_blocks = self.merger.merge(
            font_blocks=font_results,
            layout_detections=None,
            tables=tables,
        )

        # Convert to Markdown
        markdown = blocks_to_markdown(merged_blocks)

        # Post-processing steps
        markdown = self._normalize_pdf_bullets(markdown)
        markdown = self._deduplicate_pdf_blocks(markdown)
        markdown = self._fix_punctuation_spacing(markdown)

        return markdown

    # Normalization map for logical symbols and common OCR/PPTX corruption
    NORMALIZATION_MAP = {
        "": "¬",
        "": "∨",
        "": "∧",
        "": "→",
        "": "⇒",
        "": "⇔",
        "": "≡",
        "": "⊕",
        "": "↔",
        "": "∴",
        "": "∵",
        "": "∀",
        "": "∃",
        "": "∈",
        "": "∉",
        "": "∪",
        "": "∩",
        "": "⊂",
        "": "⊆",
        "": "⊋",
        "": "⊇",
        "": "∅",
        "": "√",
        "": "∞",
        "": "≠",
        "": "≈",
        "": "≤",
        "": "≥",
    }

    _PPTX_NOISE_EXACT = {
        "run",
        "animation",
        "no value",
        "reference value",
    }

    _PPTX_LABEL_HEADINGS = {
        "objective",
        "question",
        "example",
        "examples",
        "note",
        "tip",
    }

    def _normalize_text(self, text: str) -> str:
        """Apply logical symbol mapping and general text cleanup."""
        if not text:
            return text
            
        # 1. Map legacy symbols to Unicode
        for old, new in self.NORMALIZATION_MAP.items():
            text = text.replace(old, new)
            
        # 2. Fix capitalization in specific terms like "DEFINITION"
        # Only fix if it's the whole word and not an acronym
        text = re.sub(r'\bDEFINITION\b', 'Definition', text)
        text = re.sub(r'\bOBJECTIVES\b', 'Objectives', text)
        
        # 3. Detect and fix joined words (CamelCase artifacts that should have space)
        # e.g., "PeradabanAcuan" -> "Peradaban Acuan"
        # STRICTOR: Only split if the first word is 10+ chars and second 4+ chars
        # OR if it's a known non-code context. (Prevents mangling 'GeometricObject', 'printCircle', etc.)
        text = re.sub(r'([a-z]{10,})([A-Z][a-z]{4,})', r'\1 \2', text)
        
        return text

    def _normalize_slide_title_for_compare(self, text: str) -> str:
        """Normalize slide titles so continuation variants compare equal."""
        normalized = (text or "").strip().lower()
        normalized = re.sub(r"\s*[\(\[,;:-]?\s*cont\.?(?:inued)?[\)\]]?\s*$", "", normalized)
        normalized = re.sub(r"\s+", " ", normalized)
        return normalized

    def _is_probable_pptx_noise(self, text: str) -> bool:
        """Return True for common low-value PPTX artifacts like demo buttons and animation labels."""
        stripped = (text or "").strip()
        lowered = stripped.lower()
        if not stripped:
            return True
        if lowered in self._PPTX_NOISE_EXACT:
            return True
        if re.match(r"^[A-Z][A-Za-z0-9]*(?:[A-Z][A-Za-z0-9]*)+$", stripped) and len(stripped) <= 40:
            return True
        return False

    def _is_instructional_label(self, text: str) -> bool:
        """Detect short label-like headings that should render as bold body text."""
        stripped = (text or "").strip()
        label = stripped.rstrip(":").strip().lower()
        return label in self._PPTX_LABEL_HEADINGS

    def _is_probable_heading(self, text: str) -> bool:
        """Guard heading promotion so regular slide sentences are less likely to become headings."""
        stripped = (text or "").strip()
        if not stripped:
            return False
        if len(stripped) > 90:
            return False
        if stripped[-1] in ".!?":
            return False
        if self._looks_like_code(stripped):
            return False
        if re.search(r"\b(?:you|we|they|he|she|it)\b", stripped.lower()) and len(stripped.split()) >= 4:
            return False

        words = stripped.split()
        if len(words) > 9:
            return False

        return True

    def _has_substantive_slide_content(self, slide_blocks: List[Dict[str, object]]) -> bool:
        """Return True when a slide has real content beyond a bare title or demo artifacts."""
        substantive = 0
        for block in slide_blocks:
            text = str(block.get("text", "")).strip()
            btype = str(block.get("type", "body"))
            if not text or self._is_probable_pptx_noise(text):
                continue
            if btype == "h1":
                continue
            substantive += 1
        return substantive > 0

    def _is_trace_caption(self, title: str, text: str, block_type: str) -> bool:
        """Drop short animation/diagram captions on trace/demo slides."""
        title_low = (title or "").lower()
        stripped = (text or "").strip()
        if "trace" not in title_low:
            return False
        if block_type != "body":
            return False
        if self._looks_like_code(stripped):
            return False
        if len(stripped.split()) > 6:
            return False
        caption_terms = {
            "declare", "create", "assign", "change", "reference",
            "value", "circle", "mycircle", "yourcircle",
        }
        words = {w.lower().strip(".,:;()") for w in stripped.split()}
        return bool(words & caption_terms)

    def _clean_code_text(self, text: str) -> str:
        """Remove slide numbering prefixes that should not appear inside code samples."""
        cleaned = (text or "").rstrip()
        if re.match(r"^\s*\d+\.\s+(?:public|private|protected|class|interface|enum)\b", cleaned):
            cleaned = re.sub(r"^\s*\d+\.\s+", "", cleaned)
        return cleaned

    def _split_mixed_code_and_prose(self, text: str) -> Optional[tuple]:
        """
        Split a single PPTX paragraph that contains code followed by explanatory prose.
        Example:
        "Circle[] arr = new Circle[10]; An array of objects is ..."
        """
        stripped = (text or "").strip()
        if len(stripped) < 40:
            return None

        prose_starts = (
            "An ", "A ", "The ", "For ", "If ", "This ", "That ", "These ", "Those ",
            "As ", "It ", "More ", "Java ",
        )

        matches = list(re.finditer(r";\s+", stripped))
        for match in reversed(matches):
            code_part = stripped[:match.end() - 1].strip()
            prose_part = stripped[match.end():].strip()
            if not prose_part or not prose_part.startswith(prose_starts):
                continue
            if not self._looks_like_code(code_part):
                continue
            if self._looks_like_code(prose_part):
                continue
            return self._clean_code_text(code_part), prose_part

        return None

    def _postprocess_pptx_markdown(self, markdown: str) -> str:
        """Final PPTX-specific cleanup after block emission."""
        lines = markdown.splitlines()
        cleaned_lines: List[str] = []
        i = 0

        while i < len(lines):
            line = lines[i]

            # Clean the main title by dropping loud all-caps course prefixes before "Topic N".
            if i == 0 and line.startswith("# "):
                title = line[2:].strip()
                title = re.sub(
                    r"^(?:[A-Z][A-Z\-&/]+(?:\s+[A-Z][A-Z\-&/]+)+)\s+(Topic\s+\d+.*)$",
                    r"\1",
                    title,
                )
                line = f"# {title}"

            if re.fullmatch(r"(TIP|NOTE|Question):?", line.strip(), flags=re.IGNORECASE):
                label = line.strip().rstrip(":").upper()
                cleaned_lines.append(f"**{label}:**")
                i += 1
                continue

            # Remove consecutive duplicate fenced code blocks, common on trace-animation slides.
            if line.startswith("```"):
                block_start = i
                block_lines = [line]
                i += 1
                while i < len(lines):
                    block_lines.append(lines[i])
                    if lines[i].startswith("```"):
                        i += 1
                        break
                    i += 1

                block_text = "\n".join(block_lines).strip()
                prev_block = "\n".join(cleaned_lines[-len(block_lines):]).strip() if len(cleaned_lines) >= len(block_lines) else None
                if prev_block == block_text:
                    continue

                cleaned_lines.extend(block_lines)
                continue

            cleaned_lines.append(line)
            i += 1

        result = "\n".join(cleaned_lines)
        result = re.sub(r"\n{3,}", "\n\n", result).strip() + "\n"
        return result

    def _fix_punctuation_spacing(self, text: str) -> str:
        """
        Fix missing spaces after punctuation that is a common PDF extraction artefact.
        Only inserts spaces when a punctuation character is immediately followed by a
        letter or digit, but avoids common abbreviations like "v1.0", "U.S.", or dates.
        """
        import re
        # List of common abbreviations to protect (ending with a dot)
        PROTECT_PATTERNS = r'\b(v|vs|eg|ie|dr|mr|mrs|ms|prof|st|u\.s|p\.m|a\.m|p\.s)\b'
        
        # Insert a space after . , ; : when followed by a letter/digit,
        # but NOT if it looks like a version number (digit.digit) or protected abbr.
        # Negative lookbehind for common abbreviations and negative lookahead for digits (to preserve 1.0)
        text = re.sub(r'(?<![A-Z])([,;:!])([A-Za-zÀ-ž])', r'\1 \2', text)
        # For dots, be more careful: only if followed by space or end-of-sentence pattern
        text = re.sub(r'(?<![A-Z0-9\.])(\.)([A-ZÀ-ž][a-z])', r'\1 \2', text)
        return text

    def _normalize_pdf_bullets(self, text: str) -> str:
        """
        Normalize non-standard bullet characters used in lecture PDFs
        (e.g. Ø, q, n, v as line-start decorators) into proper Markdown list markers.
        """
        import re
        # These chars appear as bullet stand-ins at the start of lines
        PDF_BULLET_CHARS = r'^[ØqnvlhÂ§ø·]\s+'
        lines = text.split('\n')
        normalized = []
        for line in lines:
            stripped = line.strip()
            if re.match(PDF_BULLET_CHARS, stripped) and len(stripped) > 2:
                # Convert to a proper list item, preserving indent
                indent = len(line) - len(line.lstrip())
                content = re.sub(PDF_BULLET_CHARS, '', stripped).strip()
                normalized.append(' ' * indent + '- ' + content)
            else:
                normalized.append(line)
        return '\n'.join(normalized)

    def _deduplicate_pdf_blocks(self, markdown: str) -> str:
        """
        Remove plain-text blocks that are near-duplicates of table content.
        After pdfplumber extracts tables, the font extractor also reads those
        same characters as body text, causing each table to appear twice.
        Strategy: scan for lines directly after a table that share ≥75% of
        their words with the table rows above.
        """
        import re
        lines = markdown.split('\n')
        result = []
        # Collect table cell words for deduplication window
        table_words: set = set()
        in_table = False
        table_end_idx = -1

        for i, line in enumerate(lines):
            stripped = line.strip()

            # Track table regions
            if stripped.startswith('|'):
                in_table = True
                table_end_idx = i
                # Accumulate all words from table cells
                cells = re.split(r'\s*\|\s*', stripped)
                for cell in cells:
                    table_words.update(cell.lower().split())
                result.append(line)
                continue

            # Reset table word window after 3 non-table lines
            if in_table and i > table_end_idx + 3:
                in_table = False
                table_words = set()

            # Check candidate duplicate line (body text after a table)
            if in_table and stripped and not stripped.startswith('#') and not stripped.startswith('-'):
                line_words = set(stripped.lower().split())
                if len(line_words) >= 4 and table_words:
                    overlap = len(line_words & table_words) / len(line_words)
                    if overlap >= 0.75:
                        logger.debug(f"Dedup: skipping line with {overlap:.0%} table overlap: {stripped[:60]}")
                        continue  # Drop the duplicate

            result.append(line)

        return '\n'.join(result)

    def _extract_tables_from_pdf(self, pdf_path: str) -> list:
        """
        Extract tables from PDF and convert to markdown format.
        
        Returns a list of per-page table data:
        [
            [  # Page 1
                {"y_position": 100, "markdown": "| Header | ..."},
                ...
            ],
            [],  # Page 2 (no tables)
            ...
        ]
        """
        import pdfplumber
        
        all_tables = []
        
        try:
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    page_tables = []
                    
                    # Extract tables with their bounding boxes
                    raw_tables = page.extract_tables()
                    table_settings = page.find_tables()
                    
                    if not raw_tables:
                        all_tables.append([])
                        continue
                    
                    for idx, table in enumerate(raw_tables):
                        if not table:
                            continue
                        
                        # Convert table to markdown format
                        markdown_table = self._table_to_markdown(table)
                        
                        # Get table Y position from table settings if available
                        y_position = 0
                        if idx < len(table_settings):
                            # Get the top of the table bounding box
                            table_rect = table_settings[idx].bbox  # (x0, top, x1, bottom)
                            if table_rect:
                                y_position = table_rect[1]  # top position
                        
                        page_tables.append({
                            "y_position": y_position,
                            "markdown": markdown_table,
                        })
                    
                    all_tables.append(page_tables)
        except Exception as e:
            logger.warning(f"Table extraction failed: {e}")
            # Return empty list on failure; document extraction continues
            return []
        
        return all_tables

    def _table_to_markdown(self, table: list) -> str:
        """
        Convert a pdfplumber table (list of rows) to markdown table format.

        Args:
            table: List of rows, where each row is a list of cell contents

        Returns:
            Markdown table string
        """
        if not table or not table[0]:
            return ""

        header_row = table[0]
        markdown_lines = []

        header_cells = [str(cell or "").strip() for cell in header_row]
        markdown_lines.append("| " + " | ".join(header_cells) + " |")

        separator_cells = ["-" * max(3, len(cell)) for cell in header_cells]
        markdown_lines.append("| " + " | ".join(separator_cells) + " |")

        for row in table[1:]:
            cells = [str(cell or "").strip() for cell in row]
            markdown_lines.append("| " + " | ".join(cells) + " |")

        return "\n".join(markdown_lines)

    # Monospace fonts indicate code blocks
    MONO_FONT_KEYWORDS = ("mono", "courier", "consolas", "lucida console",
                          "inconsolata", "source code", "fira code", "jetbrains")

    def _is_monospace(self, font_name: str) -> bool:
        """Return True if the font name indicates a monospace/code font."""
        name_lower = (font_name or "").lower()
        return any(kw in name_lower for kw in self.MONO_FONT_KEYWORDS)

    # Shapes whose full text content we should skip entirely
    _SKIP_SHAPE_PATTERNS = (
        "faculty of", "department of", "university", "room no",
        "universiti", "jabatan", "fakulti",  # Malaysian university metadata
    )

    # Slides whose titles indicate they are low-value and should be skipped
    _SKIP_SLIDE_TITLES = (
        "outline", "table of contents", "learning outcomes", "objectives",
        "introduction", "summary", "conclusion", "thank you", "any questions",
        "recap", "revisions", "references", "bibliography"
    )

    def _is_metadata_shape(self, text: str, slide_num: int) -> bool:
        """Return True if this shape appears to be institutional metadata on the cover slide."""
        if slide_num != 1:
            return False
        low = text.lower()
        return any(kw in low for kw in self._SKIP_SHAPE_PATTERNS)

    def _process_pptx(self, pptx_path: str) -> str:
        """Process a PPTX file using shape-level font extraction."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for PPTX processing. "
                "Install: pip install python-pptx"
            )

        logger.info(f"Processing PPTX: {pptx_path}")
        prs = Presentation(pptx_path)
        slide_width = prs.slide_width or 1
        slide_height = prs.slide_height or 1

        # Compute the presentation's dominant body font size so we can
        # calibrate heading thresholds relative to it rather than absolutely.
        all_run_sizes = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            all_run_sizes.append(run.font.size.pt)

        if all_run_sizes:
            all_run_sizes.sort()
            median_size = all_run_sizes[len(all_run_sizes) // 2]
        else:
            median_size = 18.0

        # Heading thresholds are relative to median body size.
        # A true heading should be meaningfully larger than body text.
        h1_threshold = max(median_size * 1.6, 28.0)
        h2_threshold = max(median_size * 1.3, 22.0)
        h3_threshold = max(median_size * 1.1, 16.0)
        logger.debug(
            f"PPTX median body size: {median_size:.1f}pt → "
            f"h1≥{h1_threshold:.1f}, h2≥{h2_threshold:.1f}, h3≥{h3_threshold:.1f}"
        )

        md_parts = []
        BULLET_CHARS = set("•‣◦⁃∙‐‑–—►▪▸➤➢")
        last_slide_title = ""

        # Tracking content for cross-slide deduplication
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_blocks = []
            current_slide_title = ""

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                # Determine shape-level role
                shape_role = "body"
                # shape.top / shape.height can be None on some malformed slides
                _top = shape.top
                _height = shape.height
                shape_top_frac = (_top / slide_height) if (slide_height and _top is not None) else 1.0

                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.idx
                    # idx 0 = TITLE, idx 1 = BODY/CONTENT in standard layouts
                    if ph_type == 0:
                        shape_role = "title"
                    else:
                        shape_role = "body"
                else:
                    # Non-placeholder text box: use position heuristic
                    # Top 15% of slide and relatively short shape → likely a title
                    _h_frac = (_height / slide_height) if (slide_height and _height is not None) else 1.0
                    if shape_top_frac < 0.15 and _h_frac < 0.30:
                        shape_role = "title"

                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    # Fix PowerPoint in-shape line breaks (\x0b = vertical tab)
                    raw_text = paragraph.text.replace("\x0b", " ").rstrip()
                    # 1.2 Normalize text (symbols, capitalization)
                    raw_text = self._normalize_text(raw_text).rstrip()
                    if not raw_text.strip():
                        continue

                    # 1.3 Skip bare slide-number lines (short pure-digit strings)
                    if raw_text.strip().isdigit() and len(raw_text.strip()) <= 3:
                        logger.debug(f"Skip slide number literal: '{raw_text.strip()}'")
                        continue

                    # 1.6 Skip URLs
                    if raw_text.strip().startswith(("http://", "https://")):
                        logger.debug(f"Skip URL: '{raw_text.strip()[:60]}'")
                        continue

                    # 1.6 Skip institutional metadata on cover slide
                    if self._is_metadata_shape(raw_text.strip(), slide_num):
                        logger.debug(f"Skip metadata shape on slide 1: '{raw_text.strip()[:60]}'")
                        continue

                    text = raw_text
                    level = paragraph.level  # indent level 0–8

                    # Collect run-level font info
                    max_size = 0.0
                    is_bold = False
                    is_code = False
                    for run in paragraph.runs:
                        if run.font.size:
                            size_pt = run.font.size.pt
                            max_size = max(max_size, size_pt)
                        if run.font.bold:
                            is_bold = True
                        if self._is_monospace(run.font.name or ""):
                            is_code = True

                    # Determine block type
                    # Long text (>120 chars) is always body regardless of size.
                    is_long_text = len(text.strip()) > 120

                    # Title placeholder: ONLY the first paragraph gets h1.
                    # Subsequent paragraphs in the same shape are sub-content (body/list).
                    is_title_first_para = (shape_role == "title" and level == 0 and para_idx == 0)

                    if self._is_probable_pptx_noise(text):
                        logger.debug(f"Skip low-value PPTX artifact: '{text.strip()[:60]}'")
                        continue

                    mixed_code_prose = self._split_mixed_code_and_prose(text)
                    if mixed_code_prose:
                        code_text, prose_text = mixed_code_prose
                        slide_blocks.append({
                            "type": "code",
                            "text": code_text,
                            "indent": level,
                        })
                        slide_blocks.append({
                            "type": "body",
                            "text": prose_text.lstrip(),
                            "indent": level,
                        })
                        continue

                    if is_code or self._looks_like_code(text):
                        block_type = "code"
                    elif is_title_first_para:
                        block_type = "h1"
                        text = text.lstrip()
                    elif not is_long_text and max_size >= h1_threshold and self._is_probable_heading(text):
                        block_type = "h1"
                        text = text.lstrip()
                    elif not is_long_text and max_size >= h2_threshold and self._is_probable_heading(text):
                        block_type = "h2"
                        text = text.lstrip()
                    elif not is_long_text and max_size >= h3_threshold and is_bold and self._is_probable_heading(text):
                        block_type = "h3"
                        text = text.lstrip()
                    elif level > 0:
                        block_type = "list"
                        text = text.lstrip()
                    else:
                        # Check for common list text patterns used in slides
                        first_char = text.lstrip()[0] if text.lstrip() else ""
                        if first_char in BULLET_CHARS or text.lstrip().startswith("- "):
                            # Strip the bullet char so the list marker isn't duplicated
                            if first_char in BULLET_CHARS:
                                text = text.lstrip()[1:].strip()
                            else:
                                text = text.lstrip()[2:].strip()
                            block_type = "list"
                        else:
                            block_type = "body"
                            text = text.lstrip()

                    # Capture the title of the slide (first paragraph of title placeholder)
                    if is_title_first_para:
                        current_slide_title = text.strip()

                    slide_blocks.append({
                        "type": block_type,
                        "text": text,
                        "indent": level,
                    })

            # --- Slide-level Filtering ---
            # 1. Skip if the title is in the skip list and there's very little content
            title_lower = current_slide_title.lower().strip()
            if any(kw in title_lower for kw in self._SKIP_SLIDE_TITLES) and len(slide_blocks) <= 5:
                logger.debug(f"Skip low-value slide {slide_num}: '{current_slide_title}'")
                continue

            if slide_blocks and not self._has_substantive_slide_content(slide_blocks):
                logger.debug(f"Skip non-substantive PPTX slide {slide_num}: '{current_slide_title}'")
                continue

            # 2. Cross-slide title deduplication
            # If this slide's title is identical to the previous slide's title,
            # it's a continuation slide. Remove the title block to prevent duplication.
            is_continuation = False
            normalized_current_title = self._normalize_slide_title_for_compare(current_slide_title)
            if normalized_current_title and normalized_current_title == last_slide_title:
                is_continuation = True
                # Remove the h1 block for this slide
                slide_blocks = [b for b in slide_blocks if b["type"] != "h1"]
                logger.debug(f"Slide {slide_num} is a continuation of '{current_slide_title}'")
            
            last_slide_title = normalized_current_title

            if slide_blocks:
                slide_blocks = [
                    b for b in slide_blocks
                    if not self._is_trace_caption(current_slide_title, str(b.get("text", "")), str(b.get("type", "body")))
                ]

            if not slide_blocks:
                continue

            # --- Heading de-inflation ---
            # Count h1s beyond the first one — multiple h1s per slide means the
            # threshold still fired wrongly. Demote the excess to body/list.
            h1_blocks = [b for b in slide_blocks if b["type"] == "h1"]
            if len(h1_blocks) > 1:
                seen_h1 = False
                for b in slide_blocks:
                    if b["type"] == "h1":
                        if seen_h1:
                            # Keep as heading only if it is short (≤8 words) — genuine sub-headings
                            # A long h1 or a code-like token is clearly body text.
                            word_count = len(b["text"].split())
                            b["type"] = "h2" if word_count <= 8 else "body"
                        else:
                            seen_h1 = True

            # Demote h2/h3 blocks if they dominate (>35% of total blocks)
            h23_count = sum(1 for b in slide_blocks if b["type"] in ("h2", "h3"))
            total = len(slide_blocks)
            if total > 0 and h23_count / total > 0.35:
                logger.debug(f"Slide {slide_num}: h2/h3 inflation ({h23_count}/{total}), demoting to body")
                for b in slide_blocks:
                    if b["type"] in ("h2", "h3"):
                        b["type"] = "body"

            # Emit markdown for this slide's blocks
            in_code_block = False
            for i, block in enumerate(slide_blocks):
                btype = block["type"]
                text = block["text"]
                indent = block.get("indent", 0)

                if btype == "code":
                    text = self._clean_code_text(text)
                    if not in_code_block:
                        # Determine language hint
                        text_lower = text.lower()
                        lang = ""
                        if any(kw in text_lower for kw in ("public", "class", "void", "static", "println", "system.out")):
                            lang = "java"
                        elif any(kw in text_lower for kw in ("def ", "import ", "print(", "if __name__")):
                            lang = "python"
                        elif any(c in text for c in ("→", "∨", "∧", "¬", "↔", "≡")):
                            lang = "logic"
                        
                        md_parts.append(f"```{lang}")
                        in_code_block = True
                    prefix = "    " * indent
                    md_parts.append(f"{prefix}{text}")
                else:
                    if in_code_block:
                        md_parts.append("```")
                        md_parts.append("")
                        in_code_block = False

                    if btype.startswith("h"):
                        if self._is_instructional_label(text):
                            md_parts.append(f"**{text.rstrip(':')}:**")
                            md_parts.append("")
                            continue
                        hnum = int(btype[1])
                        md_parts.append(f"{'#' * hnum} {text}")
                        md_parts.append("")
                    elif btype == "list":
                        prefix = "  " * indent
                        md_parts.append(f"{prefix}- {text}")
                    elif btype == "ordered_list":
                        prefix = "  " * indent
                        md_parts.append(f"{prefix}1. {text}")
                    else:
                        md_parts.append(text)
                        md_parts.append("")

            if in_code_block:
                md_parts.append("```")
                md_parts.append("")

        markdown = "\n".join(md_parts).strip() + "\n"
        return self._postprocess_pptx_markdown(markdown)

    def _looks_like_code(self, text: str) -> bool:
        """Return True if text appears to be code even without monospace metadata."""
        if not text or len(text) < 10 or len(text) > 500:
            return False
            
        # Ignore lines that look like markdown headers
        if text.startswith("#"):
            return False
        
        # 1. Check for strong structural code characters (must have multiple types)
        # We look for combinations like (); or {} or []
        has_brackets = "(" in text and ")" in text
        has_braces = "{" in text and "}" in text
        has_terminate = ";" in text
        has_assignment = "=" in text and not text.startswith("=")
        
        # Require serious syntax indicators
        if (has_braces and has_terminate) or (has_brackets and has_terminate and has_assignment):
            return True

        # 2. Check for common programming keywords with strict boundary enforcement
        import re
        # Specialized keywords that are rarely used in plain lecture text without code
        keywords = ("public", "private", "class", "void", "static", "System.out", "println", "import", "def", "return", "function", "const", "let")
        kw_pattern = r'\b(' + '|'.join(keywords) + r')\b'
        matches = re.findall(kw_pattern, text)
        
        # If we see keywords like 'class' or 'public' + syntax characters, it's code
        if len(matches) >= 2 and (has_brackets or has_braces or has_terminate):
            return True

        # 3. Logic formulas (common in discrete math/logic courses)
        logic_chars = ("→", "∨", "∧", "¬", "↔", "⊕", "≡", "⇒", "⇔")
        if any(c in text for c in logic_chars) and len(text) < 150:
            # Check for pattern like "p ∧ q"
            # EXCLUDE if it contains common English/Malay words that indicate a sentence
            common_words = {"denotes", "that", "the", "has", "it", "meaning", "is", "dalam", "yang", "dan", "untuk"}
            text_lower = text.lower()
            if any(word in text_lower.split() for word in common_words):
                return False
                
            if re.search(r'[pqr]\s*[∧∨→¬]', text) or re.search(r'[¬∧∨→]\s*[pqr]', text):
                return True

        return False

    # ── chunk size for AI polish (characters) ──
    # Reduced to 1500 to minimize "stream failed" errors with slow reasoning models
    _POLISH_CHUNK_SIZE = 1500

    def _ai_polish(self, markdown: str, progress_callback: Optional[Callable[[int], None]] = None) -> str:
        """Perform a final formatting-only cleanup pass using AIClient.
        
        Splits the input into manageable chunks to prevent quality degradation
        from long contexts, then reassembles the polished chunks.
        """
        if not self.gemini_api_key or not markdown:
            return markdown

        try:
            from app.processing.ai_client import AIClient
            # Create a client which will automatically use the 3-tier fallback system
            client = AIClient() 
            
            # If a specific key was provided to the pipeline, we can override Tier 1
            if self.gemini_api_key:
                client.tiers[0].api_key = self.gemini_api_key
                client.tiers[0].model_name = self.gemini_model
                client._init_gemini_tier(client.tiers[0])

            chunks = self._split_into_chunks(markdown)
            num_chunks = len(chunks)
            logger.info(f"Refining output with {self.gemini_model} ({num_chunks} chunk(s), sequential-ish)...")

            # ── Parallel chunk processing ──
            from concurrent.futures import ThreadPoolExecutor, as_completed

            # We need an event loop for async generate_text
            import asyncio
            
            # Ensure debug directory exists for streaming chunks
            chunk_debug_dir = Path("scripts/ProcessingAlgorithmTest/output/debug_chunks")
            chunk_debug_dir.mkdir(parents=True, exist_ok=True)
            
            def _run_polish_async(idx, chunk, is_first):
                return asyncio.run(self._polish_chunk(client, idx, chunk, is_first_chunk=is_first, debug_dir=chunk_debug_dir))

            polished_chunks = [None] * num_chunks
            completed_chunks = 0

            # Reduced max_workers to 2 to prevent overwhelming reasoning models
            with ThreadPoolExecutor(max_workers=min(num_chunks, 2)) as executor:
                futures = {
                    executor.submit(_run_polish_async, i, chunk, (i == 0)): i
                    for i, chunk in enumerate(chunks)
                }
                for future in as_completed(futures):
                    idx = futures[future]
                    try:
                        result = future.result()
                        completed_chunks += 1
                        
                        # Update progress: AI starts at 30%, ends at 95%
                        if progress_callback:
                            ai_progress = 30 + int((completed_chunks / num_chunks) * 65)
                            progress_callback(ai_progress)

                        if result:
                            polished_chunks[idx] = result
                            logger.info(f"  ✓ Chunk {idx + 1}/{num_chunks} done")
                        else:
                            polished_chunks[idx] = chunks[idx]
                            logger.warning(f"  ⚠ Chunk {idx + 1}/{num_chunks} returned empty")
                    except Exception as e:
                        polished_chunks[idx] = chunks[idx]
                        logger.error(f"  ❌ Chunk {idx + 1}/{num_chunks} failed: {e}")

            result = "\n\n".join(polished_chunks).strip()

            # ── Post-processing: enforce structural rules programmatically ──
            import re
            lines = result.split("\n")
            cleaned = []
            seen_h1 = False
            for line in lines:
                # Enforce only ONE H1
                if re.match(r'^#\s+', line) and not re.match(r'^##', line):
                    if seen_h1:
                        # Demote to H2
                        line = "#" + line
                    else:
                        seen_h1 = True
                # Remove "End of Topic/Chapter" lines
                if re.match(r'^#{1,3}\s+(End of|end of)\s+(Topic|Chapter|Lecture)', line, re.IGNORECASE):
                    continue
                cleaned.append(line)

            return "\n".join(cleaned).strip()

        except Exception as e:
            logger.warning(f"AI Polish pass failed: {e}. Returning raw markdown.")

        return markdown

    def _split_into_chunks(self, markdown: str) -> List[str]:
        """Split markdown into chunks at heading boundaries to avoid mid-paragraph splits."""
        import re
        lines = markdown.split("\n")
        
        chunks: List[str] = []
        current_chunk: List[str] = []
        current_size = 0

        for line in lines:
            current_chunk.append(line)
            current_size += len(line) + 1  # +1 for newline

            # Split at heading boundaries when chunk is large enough
            if current_size >= self._POLISH_CHUNK_SIZE and re.match(r'^#{1,3}\s', line):
                # The heading line starts a new chunk
                heading_line = current_chunk.pop()
                if current_chunk:
                    chunks.append("\n".join(current_chunk))
                current_chunk = [heading_line]
                current_size = len(heading_line) + 1

        if current_chunk:
            chunks.append("\n".join(current_chunk))

        return chunks if chunks else [markdown]

    async def _polish_chunk(self, client, chunk_idx: int, chunk: str, is_first_chunk: bool = False, debug_dir: Optional[Path] = None) -> str:
        """Polish a single chunk of markdown using the AI model with streaming."""
        title_instruction = ""
        if is_first_chunk:
            title_instruction = """TITLE RULE: The first line of your output MUST be a single H1 heading (# ) with the 
EXACT topic title from the slides (e.g., "# Topic 3: Inheritance"). 
Remove university names, course codes, and lecturer names.
"""
        else:
            title_instruction = """HEADING RULE: Do NOT use H1 (# ) headings. Use only H2 (## ) or H3 (### ).
"""

        prompt = f"""Task: Clean and format the following lecture notes into clean Markdown.

CRITICAL RULES:
1. START WITH THE MARKER ===START===
2. END WITH THE MARKER ===END===
3. NO PREAMBLE/INTRO: Output ONLY the markdown between the markers.
4. NO REASONING: Do not talk to yourself, do not plan, do not list rules.
5. USE EXACT WORDS: Never rephrase or summarize.
6. ONE H1: Only the main title is #. Subtitles are ##, ###.
7. CODE BLOCKS: Use ```java only for actual code.

{title_instruction}

INPUT TO PROCESS:
{chunk}

POLISHED MARKDOWN:
===START===
"""
        try:
            full_text = ""
            debug_file = None
            if debug_dir:
                debug_file = debug_dir / f"chunk_{chunk_idx}.md"
                # Clear/create the file
                with open(debug_file, "w", encoding="utf-8") as f:
                    f.write(f"--- CHUNK {chunk_idx} START ---\n\n")

            # Use a slightly higher max_tokens to ensure we don't cut off the content
            async for text_segment in client.stream_text(prompt, max_tokens=3000):
                full_text += text_segment
                if debug_file:
                    with open(debug_file, "a", encoding="utf-8") as f:
                        f.write(text_segment)
                        f.flush()

            if full_text and not full_text.startswith("["): # Check for provider errors
                # Clean up the output
                content = full_text.strip()
                
                # Split by markers
                if "===START===" in content:
                    content = content.split("===START===")[-1]
                if "===END===" in content:
                    content = content.split("===END===")[0]
                
                content = content.strip()
                
                # Aggressive line-by-line cleanup to remove leaked reasoning artifacts
                lines = content.split("\n")
                cleaned_lines = []
                REASONING_PATTERNS = [
                    r'^\s*[\*\-]\s*Rule \d+:',
                    r'^\s*[\*\-]\s*Segment \d+:',
                    r'^\s*[\*\-]\s*Main Title \(H1\):',
                    r'^\s*[\*\-]\s*Potential H\d+',
                    r'^\s*[\*\-]\s*`# .*` \(Potential H1\)',
                    r'^\s*Correction on',
                    r'^\s*Refining the',
                    r'^\s*Let\'s verify',
                    r'^\s*Final check',
                ]
                
                import re
                for line in lines:
                    if any(re.match(p, line, re.IGNORECASE) for p in REASONING_PATTERNS):
                        continue
                    # Remove markers if they leaked into the line
                    line = line.replace("===START===", "").replace("===END===", "")
                    cleaned_lines.append(line)
                
                content = "\n".join(cleaned_lines).strip()
                
                # Final strip of any backtick wrapping (common in AI responses)
                if content.startswith("```markdown"):
                    content = content[11:].strip()
                elif content.startswith("```"):
                    content = content[3:].strip()
                if content.endswith("```"):
                    content = content[:-3].strip()

                return content
        except Exception as e:
            logger.warning(f"AI Polish chunk {chunk_idx} failed: {e}")

        return None

    def _compute_stats(self, blocks) -> Dict[str, int]:

        """Compute quality metrics for the output."""
        stats = {
            "total_blocks": len(blocks),
            "headings": 0,
            "lists": 0,
            "tables": 0,
            "body": 0,
        }
        for block in blocks:
            if block.block_type.startswith("h"):
                stats["headings"] += 1
            elif block.block_type in ("list", "ordered_list"):
                stats["lists"] += 1
            elif block.block_type == "table":
                stats["tables"] += 1
            elif block.block_type == "body":
                stats["body"] += 1
        return stats


def process_file(file_path: str, use_ai: bool = True) -> str:
    """
    Convenience function to process a file through the smart pipeline.
    
    Args:
        file_path: Path to PDF or PPTX file
        use_ai: Whether to use AI models (layout detection, table transformer)
        
    Returns:
        Clean Markdown string
    """
    pipeline = SmartPipeline(
        use_layout_detection=use_ai,
        use_table_transformer=use_ai,
    )
    return pipeline.process(file_path)
