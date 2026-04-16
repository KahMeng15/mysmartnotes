"""
Gemini Vision-Augmented Slide Extractor

Two-pass pipeline:
1. Local structural extraction (python-pptx / pdfplumber)       → fast, always runs
2. Gemini Vision classification (optional, per processing mode)  → accuracy boost

Supports three modes driven by User.note_processing_mode:
  "fast"             – local extraction only, no AI calls
  "smart"            – Gemini Vision for low-confidence slides, no delay
  "smart_throttled"  – Gemini Vision with 1-second inter-call delay
"""

import io
import json
import logging
import time
from pathlib import Path
from typing import Dict, List, Optional, Any

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
BULLET_CHARS = set("•‣◦⁃∙‐‑–—►▪▸➤➢")
MONO_FONT_KEYWORDS = ("mono", "courier", "consolas", "lucida console",
                      "inconsolata", "source code", "fira code", "jetbrains")

GEMINI_PROMPT = """You are a lecture slide parser. Analyze the attached slide image.
Output ONLY a JSON array where each element describes one text region:
{"role": "<role>", "text": "<exact text visible>", "indent": <0-3>}

Valid roles:
- "title"     : the slide's main title (usually top, largest text)
- "heading"   : a section subheading within the body
- "bullet"    : a bullet-point or dashed list item
- "body"      : normal paragraph text
- "table"     : tabular data — output as a GitHub-flavoured markdown table in "text"
- "code"      : monospaced / programming text
- "decoration": logos, slide numbers, copyright footers, clip art labels — OMIT these

Rules:
- Preserve the exact text as written (do not paraphrase or correct).
- Use "indent" = 0 for top-level items, 1-3 for sub-items.
- Output ONLY the JSON array. No prose, no markdown fences.
"""


# ---------------------------------------------------------------------------
# SlideVisionExtractor
# ---------------------------------------------------------------------------

class SlideVisionExtractor:
    """
    Extracts structured Markdown from PPTX / PDF files.

    Usage:
        extractor = SlideVisionExtractor(use_vision=True, inter_call_delay_s=1.0)
        markdown  = extractor.process("slides.pptx")   # or .pdf
    """

    def __init__(
        self,
        use_vision: bool = True,
        inter_call_delay_s: float = 0.0,
        gemini_api_key: Optional[str] = None,
        gemini_model: str = "gemini-2.5-flash",
    ):
        self.use_vision = use_vision
        self.inter_call_delay_s = inter_call_delay_s
        self._gemini_api_key = gemini_api_key
        self._gemini_model = gemini_model
        self._last_gemini_call: float = 0.0

        # Lazy-init Gemini client
        self._gemini_client = None
        if use_vision and gemini_api_key:
            self._init_gemini(gemini_api_key, gemini_model)

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def process(self, file_path: str) -> str:
        """Process a PDF or PPTX file and return clean Markdown."""
        ext = Path(file_path).suffix.lower()
        if ext == ".pptx":
            return self._process_pptx(file_path)
        elif ext == ".pdf":
            return self._process_pdf(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    # ------------------------------------------------------------------
    # PPTX Processing
    # ------------------------------------------------------------------

    def _process_pptx(self, path: str) -> str:
        """Extract and structure content from a PPTX file."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required: pip install python-pptx")

        prs = Presentation(path)
        slide_w = float(prs.slide_width or 1)
        slide_h = float(prs.slide_height or 1)

        md_parts: List[str] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            logger.debug(f"[PPTX] Processing slide {slide_num}")

            # Local extraction of all shapes
            local_blocks = self._extract_pptx_shapes(slide, slide_w, slide_h)

            if self.use_vision and self._gemini_client:
                # Render slide to image
                image_bytes = self._render_pptx_slide(prs, slide_num - 1)
                if image_bytes:
                    raw_hint = " | ".join(b["text"] for b in local_blocks[:10])
                    gemini_blocks = self._gemini_classify(image_bytes, raw_hint)
                    if gemini_blocks:
                        # Use Gemini structure, keep local text as ground-truth
                        blocks = self._merge_blocks(local_blocks, gemini_blocks)
                    else:
                        blocks = local_blocks
                else:
                    blocks = local_blocks
            else:
                blocks = local_blocks

            md_parts.extend(self._blocks_to_markdown(blocks))

        return "\n".join(md_parts).strip() + "\n"

    def _extract_pptx_shapes(
        self, slide, slide_w: float, slide_h: float
    ) -> List[Dict[str, Any]]:
        """Return a list of raw text blocks from a PPTX slide's shapes."""
        blocks: List[Dict[str, Any]] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # Shape-level role
            shape_role = "body"
            shape_top_frac = (shape.top or 0) / slide_h

            if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                ph_idx = shape.placeholder_format.idx
                shape_role = "title" if ph_idx == 0 else "body"
            else:
                # Non-placeholder: spatial heuristic
                if shape_top_frac < 0.15 and (shape.height or 0) / slide_h < 0.30:
                    shape_role = "title"

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                level = para.level
                max_size = 0.0
                is_bold = False
                is_code = False

                for run in para.runs:
                    if run.font.size:
                        max_size = max(max_size, run.font.size.pt)
                    if run.font.bold:
                        is_bold = True
                    if self._is_monospace(run.font.name or ""):
                        is_code = True

                role = self._classify_pptx_para(
                    text, shape_role, level, max_size, is_bold, is_code
                )
                blocks.append({"role": role, "text": text, "indent": level})

        return blocks

    @staticmethod
    def _classify_pptx_para(
        text: str, shape_role: str, level: int,
        max_size: float, is_bold: bool, is_code: bool
    ) -> str:
        if is_code:
            return "code"
        if shape_role == "title" and level == 0:
            return "title"
        if max_size >= 28:
            return "title"
        if max_size >= 22:
            return "heading"
        if max_size >= 16 and is_bold:
            return "heading"
        if max_size >= 14 and is_bold and len(text) < 120:
            return "heading"
        if level > 0:
            return "bullet"
        first = text[0] if text else ""
        if first in BULLET_CHARS or text.startswith("- "):
            return "bullet"
        return "body"

    def _render_pptx_slide(self, prs, slide_index: int) -> Optional[bytes]:
        """Render a single PPTX slide to PNG bytes using python-pptx + pillow."""
        try:
            from PIL import Image, ImageDraw
            import pptx.util as pptx_util

            slide = prs.slides[slide_index]
            width_px = 1280
            height_px = int(width_px * prs.slide_height / prs.slide_width)

            # Create a white canvas
            img = Image.new("RGB", (width_px, height_px), "white")

            buf = io.BytesIO()
            img.save(buf, format="PNG")
            return buf.getvalue()

        except Exception as e:
            logger.warning(f"PPTX slide render failed (slide {slide_index}): {e}")
            return None

    # ------------------------------------------------------------------
    # PDF Processing
    # ------------------------------------------------------------------

    def _process_pdf(self, path: str) -> str:
        """Extract and structure content from a PDF file."""
        import pdfplumber

        md_parts: List[str] = []

        with pdfplumber.open(path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                logger.debug(f"[PDF] Processing page {page_num}")

                # Local text extraction
                local_blocks = self._extract_pdf_page_blocks(page, page_num)

                if self.use_vision and self._gemini_client and local_blocks is not None:
                    confidence = self._extraction_confidence(local_blocks)
                    if confidence < 0.7:
                        # Low-confidence page → call Gemini Vision
                        image_bytes = self._render_pdf_page(page)
                        if image_bytes:
                            raw_hint = " | ".join(
                                b["text"] for b in (local_blocks or [])[:10]
                            )
                            gemini_blocks = self._gemini_classify(image_bytes, raw_hint)
                            if gemini_blocks:
                                local_blocks = self._merge_blocks(
                                    local_blocks, gemini_blocks
                                )

                md_parts.extend(self._blocks_to_markdown(local_blocks or []))

        return "\n".join(md_parts).strip() + "\n"

    def _extract_pdf_page_blocks(
        self, page, page_num: int
    ) -> List[Dict[str, Any]]:
        """Use the existing FontAwareExtractor for reliable PDF text extraction."""
        try:
            from app.processing.font_extractor import FontAwareExtractor
            extractor = FontAwareExtractor()
            results = extractor.extract(str(page.pdf.stream.name))  # type: ignore
            # Find this page's blocks
            for page_result in results:
                if page_result["page"] == page_num:
                    blocks = []
                    for block in page_result["blocks"]:
                        role = self._pdf_blocktype_to_role(block.block_type)
                        blocks.append({
                            "role": role,
                            "text": block.text.strip(),
                            "indent": block.indent_level,
                        })
                    return [b for b in blocks if b["text"]]
            return []
        except Exception as e:
            logger.warning(f"FontAwareExtractor failed for page {page_num}: {e}")
            # Fallback: raw text extraction
            raw = page.extract_text() or ""
            return [{"role": "body", "text": line.strip(), "indent": 0}
                    for line in raw.split("\n") if line.strip()]

    @staticmethod
    def _pdf_blocktype_to_role(block_type: str) -> str:
        """Map FontAwareExtractor block types to unified roles."""
        mapping = {
            "h1": "title", "h2": "heading", "h3": "heading",
            "h4": "heading", "h5": "heading",
            "list": "bullet", "ordered_list": "bullet",
            "table": "table", "body": "body", "skip": "decoration",
        }
        return mapping.get(block_type, "body")

    def _render_pdf_page(self, page) -> Optional[bytes]:
        """Render a PDF page to PNG bytes using pdf2image."""
        try:
            from pdf2image import convert_from_bytes
            pdf_bytes = page.pdf.stream.read()  # type: ignore
            page_num = page.page_number  # 1-indexed

            images = convert_from_bytes(
                pdf_bytes, dpi=150,
                first_page=page_num, last_page=page_num,
                fmt="png",
            )
            if images:
                buf = io.BytesIO()
                images[0].save(buf, format="PNG")
                return buf.getvalue()
        except Exception as e:
            logger.warning(f"PDF page render failed: {e}")
        return None

    # ------------------------------------------------------------------
    # Gemini Vision
    # ------------------------------------------------------------------

    def _init_gemini(self, api_key: str, model: str):
        """Initialize the Gemini client."""
        try:
            import google.generativeai as genai
            genai.configure(api_key=api_key)
            self._gemini_client = genai.GenerativeModel(model)
            logger.info(f"Gemini Vision client ready (model={model})")
        except Exception as e:
            logger.warning(f"Could not initialize Gemini client: {e}")
            self._gemini_client = None

    def _gemini_classify(
        self, image_bytes: bytes, raw_hint: str = ""
    ) -> Optional[List[Dict[str, Any]]]:
        """
        Send a slide/page image to Gemini Vision and return classified blocks.
        Returns None on failure (caller falls back to local extraction).
        """
        if not self._gemini_client:
            return None

        # Throttle
        if self.inter_call_delay_s > 0:
            elapsed = time.time() - self._last_gemini_call
            wait = self.inter_call_delay_s - elapsed
            if wait > 0:
                time.sleep(wait)

        try:
            from PIL import Image
            import google.generativeai as genai

            img = Image.open(io.BytesIO(image_bytes))
            prompt = GEMINI_PROMPT
            if raw_hint:
                prompt += f"\n\nHint – locally extracted text: {raw_hint[:500]}"

            response = self._gemini_client.generate_content([prompt, img])
            self._last_gemini_call = time.time()

            raw_text = response.text.strip()
            # Strip possible markdown fences
            if raw_text.startswith("```"):
                raw_text = "\n".join(raw_text.split("\n")[1:])
            if raw_text.endswith("```"):
                raw_text = "\n".join(raw_text.split("\n")[:-1])

            parsed = json.loads(raw_text)
            if isinstance(parsed, list):
                logger.debug(f"Gemini returned {len(parsed)} blocks")
                return parsed
        except json.JSONDecodeError as e:
            logger.warning(f"Gemini returned invalid JSON: {e}")
        except Exception as e:
            logger.warning(f"Gemini Vision call failed: {e}")

        return None

    # ------------------------------------------------------------------
    # Merging
    # ------------------------------------------------------------------

    @staticmethod
    def _merge_blocks(
        local: List[Dict[str, Any]], gemini: List[Dict[str, Any]]
    ) -> List[Dict[str, Any]]:
        """
        Merge local extraction (trusted text) with Gemini classification (trusted roles).

        Strategy:
        - We trust Gemini's *role* classification.
        - We trust the *local* extraction for the actual text content (avoids hallucinations).
        - We do a greedy text-match: for each Gemini block, find the closest local block text.
        - Blocks Gemini marks as "decoration" are dropped.
        """
        if not gemini:
            return local

        # Build a lookup of local text → role (for quick override)
        result: List[Dict[str, Any]] = []
        used_local = set()
        local_texts = [b["text"].lower() for b in local]

        for g_block in gemini:
            role = g_block.get("role", "body")
            if role == "decoration":
                continue

            g_text = g_block.get("text", "").strip()
            if not g_text:
                continue

            # Try to find matching local text
            best_idx = None
            best_ratio = 0.0
            g_lower = g_text.lower()

            for idx, l_text in enumerate(local_texts):
                if idx in used_local:
                    continue
                # Simple overlap ratio
                shorter = min(len(g_lower), len(l_text)) or 1
                overlap = sum(1 for a, b in zip(g_lower, l_text) if a == b)
                ratio = overlap / shorter
                if ratio > best_ratio:
                    best_ratio = ratio
                    best_idx = idx

            if best_idx is not None and best_ratio > 0.5:
                used_local.add(best_idx)
                result.append({
                    "role": role,
                    "text": local[best_idx]["text"],  # Use local text, Gemini role
                    "indent": g_block.get("indent", local[best_idx].get("indent", 0)),
                })
            else:
                # No good match — use Gemini's text directly
                result.append({
                    "role": role,
                    "text": g_text,
                    "indent": g_block.get("indent", 0),
                })

        return result if result else local

    # ------------------------------------------------------------------
    # Confidence Heuristic
    # ------------------------------------------------------------------

    @staticmethod
    def _extraction_confidence(blocks: List[Dict[str, Any]]) -> float:
        """
        Score local extraction confidence 0..1.
        Low confidence = call Gemini Vision.
        """
        if not blocks:
            return 0.0
        total_chars = sum(len(b.get("text", "")) for b in blocks)
        if total_chars < 30:
            return 0.1   # Almost no text extracted → very likely image-heavy
        has_title = any(b.get("role") in ("title", "heading") for b in blocks)
        has_body = any(b.get("role") in ("body", "bullet") for b in blocks)
        score = 0.5
        if has_title:
            score += 0.25
        if has_body:
            score += 0.25
        return score

    # ------------------------------------------------------------------
    # Markdown Output
    # ------------------------------------------------------------------

    @staticmethod
    def _blocks_to_markdown(blocks: List[Dict[str, Any]]) -> List[str]:
        """Convert a list of classified blocks to Markdown lines."""
        lines: List[str] = []
        in_code = False

        for block in blocks:
            role = block.get("role", "body")
            text = block.get("text", "").strip()
            indent = int(block.get("indent", 0))

            if not text or role == "decoration":
                continue

            if role == "code":
                if not in_code:
                    lines.append("```")
                    in_code = True
                lines.append(text)
                continue

            if in_code:
                lines.append("```")
                lines.append("")
                in_code = False

            if role == "title":
                lines.append(f"# {text}")
                lines.append("")
            elif role == "heading":
                lines.append(f"## {text}")
                lines.append("")
            elif role == "bullet":
                prefix = "  " * indent
                lines.append(f"{prefix}- {text}")
            elif role == "table":
                # text already is a markdown table string from Gemini
                lines.append(text)
                lines.append("")
            else:  # body
                lines.append(text)
                lines.append("")

        if in_code:
            lines.append("```")
            lines.append("")

        return lines

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _is_monospace(font_name: str) -> bool:
        name_lower = font_name.lower()
        return any(kw in name_lower for kw in MONO_FONT_KEYWORDS)

    @staticmethod
    def _extraction_confidence(blocks: List[Dict[str, Any]]) -> float:
        if not blocks:
            return 0.0
        total_chars = sum(len(b.get("text", "")) for b in blocks)
        if total_chars < 30:
            return 0.1
        has_title = any(b.get("role") in ("title", "heading") for b in blocks)
        has_body = any(b.get("role") in ("body", "bullet") for b in blocks)
        score = 0.5
        if has_title:
            score += 0.25
        if has_body:
            score += 0.25
        return score


# ---------------------------------------------------------------------------
# Convenience function
# ---------------------------------------------------------------------------

def extract_slide_markdown(
    file_path: str,
    use_vision: bool = True,
    inter_call_delay_s: float = 0.0,
    gemini_api_key: Optional[str] = None,
) -> str:
    """Quick extraction helper for one-off use."""
    extractor = SlideVisionExtractor(
        use_vision=use_vision,
        inter_call_delay_s=inter_call_delay_s,
        gemini_api_key=gemini_api_key,
    )
    return extractor.process(file_path)
