"""OCR and text extraction module"""
import json
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
import logging
from typing import List, Dict, Any, Tuple

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from app.processing.text_processor import process_extracted_text, ContentSegment

from app.processing.image_extractor import ImageExtractor

logger = logging.getLogger(__name__)


class OCRProcessor:
    """Handle text extraction from various file types with structured output"""
    
    @staticmethod
    def extract_from_pdf(file_path: str, lecture_id: int = 0, structured: bool = True, use_v2: bool = True) -> Tuple[str, List[ContentSegment], List]:
        """
        Extract text from PDF using OCR with optional structured processing
        
        Args:
            file_path: Path to PDF file
            lecture_id: Lecture ID for organizing extracted images
            structured: Whether to return structured segments
            use_v2: Whether to use enhanced v2 text processor (default True)
            
        Returns:
            Tuple of (raw_text, structured_segments, images)
        """
        try:
            logger.info(f"Extracting text from PDF: {file_path} (use_v2={use_v2})")
            images = convert_from_path(file_path)
            text = ""
            all_segments = []
            
            for page_num, image in enumerate(images):
                logger.debug(f"Processing PDF page {page_num + 1}")
                # Convert image to text using Tesseract
                page_text = pytesseract.image_to_string(image)
                text += f"\n--- Page {page_num + 1} ---\n{page_text}"
                
                # Process into structured segments if requested
                if structured:
                    segments = process_extracted_text(page_text, page_number=page_num + 1)
                    logger.debug(f"Processed with v1: {len(segments)} segments from page {page_num + 1}")
                    all_segments.extend(segments)
            
            # Extract images if lecture_id provided
            images_extracted = []
            if lecture_id > 0:
                try:
                    image_extractor = ImageExtractor(lecture_id=lecture_id)
                    images_extracted = image_extractor.extract_images_from_pdf(file_path)
                    logger.info(f"Extracted {len(images_extracted)} images from PDF")
                except Exception as e:
                    logger.warning(f"Error extracting images: {e}")
            
            return text.strip(), all_segments, images_extracted
            
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")
            raise
    
    @staticmethod
    def extract_from_pptx(file_path: str, lecture_id: int = 0, structured: bool = True, use_v2: bool = True) -> Tuple[str, List[ContentSegment], List]:
        """
        Extract text from PowerPoint presentation
        
        Args:
            file_path: Path to PPTX file
            lecture_id: Lecture ID for organizing extracted images
            structured: Whether to return structured segments
            use_v2: Whether to use enhanced v2 text processor (default True)
            
        Returns:
            Tuple of (raw_text, structured_segments, images)
        """
        if not HAS_PPTX:
            logger.error("python-pptx module not available. Install with: pip install python-pptx")
            raise ImportError("python-pptx is not installed. Please install it with: pip install python-pptx")
        try:
            logger.info(f"Extracting text from PPTX: {file_path} (use_v2={use_v2})")
            prs = Presentation(file_path)
            text = ""
            all_segments = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"\n--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                
                text += slide_text
                
                # Process into structured segments if requested
                # Process into structured segments if requested
                if structured:
                    segments = process_extracted_text(slide_text, page_number=slide_num + 1)
                    logger.debug(f"Processed with v1: {len(segments)} segments from slide {slide_num + 1}")
                    all_segments.extend(segments)
            
            return text.strip(), all_segments, []
            
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {e}")
            raise
    
    @staticmethod
    def extract_from_image(file_path: str, page_number: int = 1, structured: bool = True, use_v2: bool = True) -> Tuple[str, List[ContentSegment]]:
        """
        Extract text from image using OCR
        
        Args:
            file_path: Path to image file
            page_number: Page/image number for reference
            structured: Whether to return structured segments
            use_v2: Whether to use enhanced v2 text processor (default True)
            
        Returns:
            Tuple of (raw_text, structured_segments)
        """
        try:
            logger.info(f"Extracting text from image: {file_path} (use_v2={use_v2})")
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            
            all_segments = []
            if structured:
                all_segments = process_extracted_text(text, page_number=page_number)
                logger.debug(f"Processed with v1: {len(all_segments)} segments from image")
            
            return text.strip(), all_segments
            
        except Exception as e:
            logger.error(f"Error extracting text from image: {e}")
            raise
    
    @staticmethod
    def extract_text(file_path: str, file_type: str, lecture_id: int = 0, use_v2: bool = True) -> Dict[str, Any]:
        """
        Extract text based on file type with full structured output
        
        Args:
            file_path: Path to the file
            file_type: MIME type or file extension
            lecture_id: Lecture ID for organizing output
            use_v2: Whether to use enhanced v2 text processor (default True)
            
        Returns:
            Dict with raw_text, structured_content, images, and metadata
        """
        try:
            if "pdf" in file_type.lower():
                raw_text, segments, images = OCRProcessor.extract_from_pdf(
                    file_path, 
                    lecture_id=lecture_id, 
                    structured=True,
                    use_v2=use_v2
                )
            elif "pptx" in file_type.lower() or "presentation" in file_type.lower():
                raw_text, segments, images = OCRProcessor.extract_from_pptx(
                    file_path, 
                    lecture_id=lecture_id, 
                    structured=True,
                    use_v2=use_v2
                )
            elif "image" in file_type.lower() or file_type.lower().endswith((".png", ".jpg", ".jpeg")):
                raw_text, segments = OCRProcessor.extract_from_image(file_path, structured=True, use_v2=use_v2)
                images = []
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
            
            # Convert segments to dictionaries for JSON serialization
            structured_content = [s.to_dict() for s in segments]
            images_data = [img.to_dict() for img in images]
            
            return {
                "raw_text": raw_text,
                "structured_content": structured_content,
                "images": images_data,
                "segment_count": len(segments),
                "image_count": len(images),
                "content_types": {
                    "h1": sum(1 for s in segments if s.content_type.value == "h1"),
                    "h2": sum(1 for s in segments if s.content_type.value == "h2"),
                    "h3": sum(1 for s in segments if s.content_type.value == "h3"),
                    "body": sum(1 for s in segments if s.content_type.value == "body"),
                    "list": sum(1 for s in segments if s.content_type.value == "list"),
                },
            }
        except Exception as e:
            logger.error(f"Error extracting text: {e}")
            raise
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list:
        """
        Split text into chunks for embedding
        
        Args:
            text: Text to chunk
            chunk_size: Size of each chunk (words)
            overlap: Number of overlapping words between chunks
            
        Returns:
            List of text chunks
        """
        words = text.split()
        chunks = []
        
        for i in range(0, len(words), chunk_size - overlap):
            chunk = " ".join(words[i:i + chunk_size])
            if chunk.strip():
                chunks.append(chunk)
        
        return chunks
