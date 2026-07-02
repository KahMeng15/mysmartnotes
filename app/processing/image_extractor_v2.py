"""
Multi-format image extraction engine.

Extracts images from PDF, PPTX, and DOCX files with full positional metadata.
Filters out logos, backgrounds, and decorative elements using heuristics.
"""

import logging
import os
import uuid
from dataclasses import asdict, dataclass
from pathlib import Path

from app.processing.pipeline_knowledge import PipelineKnowledge

logger = logging.getLogger(__name__)


@dataclass
class ExtractedImage:
    id: str
    filename: str
    file_path: str
    source_file: str
    format: str  # pdf, pptx, docx, image
    page_number: int = 0
    slide_index: int = 0
    position_x: float = 0.0
    position_y: float = 0.0
    width: float = 0.0
    height: float = 0.0
    bbox: dict | None = None
    source_shape_type: str = "unknown"  # picture, chart, smartart, group, embedded, contour
    is_diagram: bool = False
    is_decorative: bool = False
    confidence: float = 0.8
    caption: str = ""
    alt_text: str = ""
    md5_hash: str = ""

    def to_dict(self):
        return asdict(self)


class ImageClassifier:
    """Classifies extracted images as diagram, photo, logo, background, or decorative."""

    MIN_IMAGE_SIZE = 150
    CORNER_MARGIN_RATIO = 0.05
    FULL_BLEED_RATIO = 0.85
    LOW_VARIANCE_THRESHOLD = 15
    EXTREME_ASPECT = 5

    def __init__(self, knowledge: PipelineKnowledge | None = None):
        self._knowledge = knowledge or PipelineKnowledge()
        self._position_cache: dict[str, list] = {}  # filename -> [(x, y, w, h), ...]

    def classify(
        self, image: ExtractedImage, page_width: float = 0, page_height: float = 0
    ) -> ExtractedImage:
        if image.width < self.MIN_IMAGE_SIZE and image.height < self.MIN_IMAGE_SIZE:
            image.is_decorative = True
            image.confidence = 0.95
            image.is_diagram = False
            return image

        if page_width and page_height:
            x_ratio = image.position_x / page_width if page_width else 0
            y_ratio = image.position_y / page_height if page_height else 0
            w_ratio = image.width / page_width if page_width else 0
            h_ratio = image.height / page_height if page_height else 0

            if w_ratio > self.FULL_BLEED_RATIO and h_ratio > self.FULL_BLEED_RATIO:
                image.is_decorative = True
                image.confidence = 0.9
                image.is_diagram = False
                return image

            if (
                x_ratio < self.CORNER_MARGIN_RATIO or x_ratio > (1 - self.CORNER_MARGIN_RATIO)
            ) and y_ratio < self.CORNER_MARGIN_RATIO:
                image.is_decorative = True
                image.confidence = 0.85
                image.is_diagram = False
                return image

        if image.width > 0 and image.height > 0:
            aspect = image.width / image.height
            if aspect > self.EXTREME_ASPECT or (1 / aspect) > self.EXTREME_ASPECT:
                image.is_decorative = True
                image.confidence = 0.8
                image.is_diagram = False
                return image

        if image.source_shape_type in ("chart", "smartart"):
            image.is_diagram = True
            image.confidence = 0.95
            image.is_decorative = False
            return image

        if image.source_shape_type == "picture":
            image.is_diagram = True
            image.confidence = 0.7
            image.is_decorative = False
            return image

        image.is_diagram = True
        image.confidence = 0.5
        image.is_decorative = False
        return image

    def check_template_repeat(self, image: ExtractedImage, source_key: str) -> ExtractedImage:
        if source_key not in self._position_cache:
            self._position_cache[source_key] = []
        pos_key = (
            round(image.position_x, -1),
            round(image.position_y, -1),
            round(image.width, -1),
            round(image.height, -1),
        )
        self._position_cache[source_key].append(pos_key)

        positions = self._position_cache[source_key]
        count = sum(1 for p in positions if p == pos_key)
        if count >= 3:
            image.is_decorative = True
            image.is_diagram = False
            image.confidence = 0.95
            return image

        if image.md5_hash and self._knowledge.should_skip_image(image.md5_hash):
            image.is_decorative = True
            image.is_diagram = False
            image.confidence = 0.98
            return image

        patterns = self._knowledge.decorative_patterns
        for pattern in patterns:
            max_w = pattern.get("max_width", 0)
            max_h = pattern.get("max_height", 0)
            if max_w and max_h and image.width <= max_w and image.height <= max_h:
                image.is_decorative = True
                image.is_diagram = False
                image.confidence = pattern.get("confidence", 0.8)
                return image

        return image


class ImageExtractorV2:
    """Extracts images from PDF, PPTX, DOCX with position metadata."""

    SUPPORTED_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}

    def __init__(
        self,
        output_base_dir: str | None = None,
        classifier: ImageClassifier | None = None,
    ):
        from app.utils.paths import LEGACY_EXTRACTED_IMAGES_DIR
        self.output_base_dir = output_base_dir or LEGACY_EXTRACTED_IMAGES_DIR
        self.classifier = classifier or ImageClassifier()

    def extract(self, file_path: str, resource_id: str = "") -> list[ExtractedImage]:
        ext = Path(file_path).suffix.lower()

        from app.utils.storage import StorageManager
        output_dir = StorageManager.get_extracted_images_dir_for_resource(
            resource_id or Path(file_path).stem
        )

        if ext == ".pdf":
            return self._extract_from_pdf(file_path, output_dir, resource_id)
        elif ext == ".pptx":
            return self._extract_from_pptx(file_path, output_dir, resource_id)
        elif ext == ".docx":
            return self._extract_from_docx(file_path, output_dir, resource_id)
        elif ext in self.SUPPORTED_IMAGE_EXTS:
            return self._extract_from_image(file_path, output_dir, resource_id)
        return []

    def _save_image_blob(self, blob: bytes, output_dir: str, prefix: str, ext: str = ".png") -> str:
        filename = f"{prefix}_{uuid.uuid4().hex[:8]}{ext}"
        filepath = os.path.join(output_dir, filename)
        with open(filepath, "wb") as f:
            f.write(blob)
        return filepath

    def _extract_from_pptx(
        self, pptx_path: str, output_dir: str, resource_id: str
    ) -> list[ExtractedImage]:
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx not installed")
            return []

        images = []
        prs = Presentation(pptx_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for slide_idx, slide in enumerate(prs.slides):
            self._extract_shapes_recursive(
                slide.shapes,
                slide_idx,
                slide_width,
                slide_height,
                output_dir,
                pptx_path,
                images,
                resource_id,
            )

        return self._deduplicate_and_classify(images, pptx_path)

    def _extract_shapes_recursive(
        self,
        shapes,
        slide_idx,
        slide_width,
        slide_height,
        output_dir,
        source_path,
        images,
        resource_id,
        depth=0,
    ):
        if depth > 10:
            return
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            return

        for shape in shapes:
            try:
                shape_type = getattr(shape, "shape_type", None)
                if shape_type == MSO_SHAPE_TYPE.GROUP:
                    self._extract_shapes_recursive(
                        shape.shapes,
                        slide_idx,
                        slide_width,
                        slide_height,
                        output_dir,
                        source_path,
                        images,
                        resource_id,
                        depth + 1,
                    )
                    continue

                if shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_id = f"img_{uuid.uuid4().hex[:8]}"
                    try:
                        image_blob = shape.image.blob
                        img_ext = (
                            "." + shape.image.content_type.split("/")[-1]
                            if "/" in (shape.image.content_type or "")
                            else ".png"
                        )
                        file_path = self._save_image_blob(
                            image_blob, output_dir, f"slide{slide_idx + 1}_{img_id}", img_ext
                        )

                        extracted = ExtractedImage(
                            id=img_id,
                            filename=os.path.basename(file_path),
                            file_path=file_path,
                            source_file=source_path,
                            format="pptx",
                            slide_index=slide_idx + 1,
                            page_number=slide_idx + 1,
                            position_x=getattr(shape, "left", 0) or 0,
                            position_y=getattr(shape, "top", 0) or 0,
                            width=getattr(shape, "width", 0) or 0,
                            height=getattr(shape, "height", 0) or 0,
                            bbox={
                                "x": getattr(shape, "left", 0),
                                "y": getattr(shape, "top", 0),
                                "w": getattr(shape, "width", 0),
                                "h": getattr(shape, "height", 0),
                            },
                            source_shape_type="picture",
                        )
                        images.append(extracted)
                    except Exception as e:
                        logger.debug(f"Failed to extract picture from slide {slide_idx + 1}: {e}")

                elif hasattr(shape, "has_chart") and shape.has_chart:
                    img_id = f"chart_{uuid.uuid4().hex[:8]}"
                    extracted = ExtractedImage(
                        id=img_id,
                        filename="",
                        file_path="",
                        source_file=source_path,
                        format="pptx",
                        slide_index=slide_idx + 1,
                        page_number=slide_idx + 1,
                        position_x=getattr(shape, "left", 0) or 0,
                        position_y=getattr(shape, "top", 0) or 0,
                        width=getattr(shape, "width", 0) or 0,
                        height=getattr(shape, "height", 0) or 0,
                        bbox={
                            "x": getattr(shape, "left", 0),
                            "y": getattr(shape, "top", 0),
                            "w": getattr(shape, "width", 0),
                            "h": getattr(shape, "height", 0),
                        },
                        source_shape_type="chart",
                        caption=shape.chart.chart_title.text
                        if hasattr(shape.chart, "chart_title")
                        else "",
                    )
                    images.append(extracted)

                elif hasattr(shape, "has_smart_art") and shape.has_smart_art:
                    img_id = f"smartart_{uuid.uuid4().hex[:8]}"
                    extracted = ExtractedImage(
                        id=img_id,
                        filename="",
                        file_path="",
                        source_file=source_path,
                        format="pptx",
                        slide_index=slide_idx + 1,
                        page_number=slide_idx + 1,
                        position_x=getattr(shape, "left", 0) or 0,
                        position_y=getattr(shape, "top", 0) or 0,
                        width=getattr(shape, "width", 0) or 0,
                        height=getattr(shape, "height", 0) or 0,
                        bbox={
                            "x": getattr(shape, "left", 0),
                            "y": getattr(shape, "top", 0),
                            "w": getattr(shape, "width", 0),
                            "h": getattr(shape, "height", 0),
                        },
                        source_shape_type="smartart",
                    )
                    images.append(extracted)

            except Exception as e:
                logger.debug(f"Error processing shape on slide {slide_idx + 1}: {e}")
                continue

    def _extract_from_pdf(
        self, pdf_path: str, output_dir: str, resource_id: str
    ) -> list[ExtractedImage]:
        images = []

        try:
            import fitz

            doc = fitz.open(pdf_path)
            for page_num, page in enumerate(doc):
                for img_index, img in enumerate(page.get_images(full=True)):
                    try:
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        img_ext = "." + base_image["ext"]

                        img_id = f"p{page_num + 1}_e{img_index}"
                        file_path = self._save_image_blob(image_bytes, output_dir, img_id, img_ext)

                        bbox = None
                        for image_info in page.get_image_info():
                            if image_info.get("xref") == xref:
                                bbox = image_info.get("bbox")
                                break

                        extracted = ExtractedImage(
                            id=img_id,
                            filename=os.path.basename(file_path),
                            file_path=file_path,
                            source_file=pdf_path,
                            format="pdf",
                            page_number=page_num + 1,
                            position_x=bbox[0] if bbox else 0,
                            position_y=bbox[1] if bbox else 0,
                            width=(bbox[2] - bbox[0]) if bbox else 0,
                            height=(bbox[3] - bbox[1]) if bbox else 0,
                            bbox={"x0": bbox[0], "y0": bbox[1], "x1": bbox[2], "y1": bbox[3]}
                            if bbox
                            else None,
                            source_shape_type="embedded",
                        )
                        images.append(extracted)
                    except Exception as e:
                        logger.debug(f"Failed to extract image from PDF page {page_num + 1}: {e}")
            doc.close()
        except ImportError:
            logger.debug("PyMuPDF (fitz) not available for PDF image extraction")
        except Exception as e:
            logger.warning(f"PDF image extraction error: {e}")

        try:
            import cv2
            import numpy as np
            import pytesseract
            from pdf2image import convert_from_path

            pil_images = convert_from_path(pdf_path)
            for page_num, pil_img in enumerate(pil_images):
                cv_img = cv2.cvtColor(np.array(pil_img), cv2.COLOR_RGB2BGR)
                gray = cv2.cvtColor(cv_img, cv2.COLOR_BGR2GRAY)
                grad = cv2.morphologyEx(
                    gray, cv2.MORPH_GRADIENT, cv2.getStructuringElement(cv2.MORPH_RECT, (2, 2))
                )
                _, binary = cv2.threshold(grad, 0, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)
                kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (50, 20))
                connected = cv2.morphologyEx(binary, cv2.MORPH_CLOSE, kernel)
                contours, _ = cv2.findContours(
                    connected, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE
                )

                height, width = cv_img.shape[:2]
                for cnt in contours:
                    x, y, w, h = cv2.boundingRect(cnt)
                    if w < 120 or h < 120:
                        continue
                    if w > width * 0.9 and h > height * 0.9:
                        continue
                    touches_side = (x < 10) or (x + w > width - 10)
                    if touches_side and h > height * 0.5:
                        continue
                    if w / h > 12 or h / w > 12:
                        continue

                    roi_gray = gray[y : y + h, x : x + w]
                    std_dev = np.std(roi_gray)
                    if std_dev < 20:
                        continue

                    try:
                        text_content = pytesseract.image_to_string(
                            roi_gray, config="--psm 6"
                        ).strip()
                        if 0 < len(text_content) < 4:
                            continue
                        if text_content.isalpha() and len(text_content) < 6:
                            continue
                    except Exception:
                        pass

                    existing_bbox = [e for e in images if e.page_number == page_num + 1]
                    is_duplicate = False
                    for existing in existing_bbox:
                        if existing.bbox:
                            ex, ey = existing.bbox.get("x0", 0), existing.bbox.get("y0", 0)
                            if abs(ex - x) < 50 and abs(ey - y) < 50:
                                is_duplicate = True
                                break
                    if is_duplicate:
                        continue

                    img_id = f"p{page_num + 1}_cv{uuid.uuid4().hex[:6]}"
                    roi = cv_img[y : y + h, x : x + w]
                    filename = f"{img_id}.jpg"
                    save_path = os.path.join(output_dir, filename)
                    cv2.imwrite(save_path, roi)

                    extracted = ExtractedImage(
                        id=img_id,
                        filename=filename,
                        file_path=save_path,
                        source_file=pdf_path,
                        format="pdf",
                        page_number=page_num + 1,
                        position_x=x,
                        position_y=y,
                        width=w,
                        height=h,
                        bbox={"x0": x, "y0": y, "x1": x + w, "y1": y + h},
                        source_shape_type="contour",
                    )
                    images.append(extracted)

        except ImportError:
            logger.debug("OpenCV/pdf2image not available for contour-based image extraction")
        except Exception as e:
            logger.warning(f"Contour image extraction error: {e}")

        return self._deduplicate_and_classify(images, pdf_path)

    def _extract_from_docx(
        self, docx_path: str, output_dir: str, resource_id: str
    ) -> list[ExtractedImage]:
        try:
            from docx import Document
        except ImportError:
            logger.error("python-docx not installed")
            return []

        images = []
        doc = Document(docx_path)
        nsmap = {
            "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
            "wp": "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        }

        for para_idx, paragraph in enumerate(doc.paragraphs):
            drawings = paragraph._element.findall(".//w:drawing", nsmap)
            if not drawings:
                drawings = paragraph._element.findall(".//wp:inline", nsmap)
            if not drawings:
                drawings = paragraph._element.findall(".//wp:anchor", nsmap)

            for drawing in drawings:
                try:
                    blip = drawing.findall(".//a:blip", nsmap)
                    if not blip:
                        continue
                    for b in blip:
                        embed_id = b.get(
                            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                        )
                        if not embed_id:
                            continue
                        rel = doc.part.related_parts.get(embed_id)
                        if not rel:
                            continue
                        image_blob = rel.blob
                        content_type = getattr(rel, "content_type", "image/png") or "image/png"
                        img_ext = (
                            "." + content_type.split("/")[-1] if "/" in content_type else ".png"
                        )

                        img_id = f"docx_p{para_idx}_{uuid.uuid4().hex[:6]}"
                        file_path = self._save_image_blob(image_blob, output_dir, img_id, img_ext)

                        extent = drawing.findall(".//wp:extent", nsmap)
                        cx = cy = 0
                        if extent:
                            cx = int(extent[0].get("cx", 0))
                            cy = int(extent[0].get("cy", 0))

                        extracted = ExtractedImage(
                            id=img_id,
                            filename=os.path.basename(file_path),
                            file_path=file_path,
                            source_file=docx_path,
                            format="docx",
                            page_number=0,
                            position_x=0,
                            position_y=para_idx * 20,
                            width=cx / 914400 * 96,
                            height=cy / 914400 * 96,
                            source_shape_type="inline",
                        )
                        images.append(extracted)
                except Exception as e:
                    logger.debug(f"Failed to extract image from DOCX paragraph {para_idx}: {e}")

        return self._deduplicate_and_classify(images, docx_path)

    def _extract_from_image(
        self, image_path: str, output_dir: str, resource_id: str
    ) -> list[ExtractedImage]:
        import shutil

        filename = os.path.basename(image_path)
        dest_path = os.path.join(output_dir, filename)
        shutil.copy2(image_path, dest_path)

        img_id = f"img_{uuid.uuid4().hex[:8]}"
        extracted = ExtractedImage(
            id=img_id,
            filename=filename,
            file_path=dest_path,
            source_file=image_path,
            format="image",
            page_number=1,
            source_shape_type="direct_image",
            is_diagram=True,
            confidence=0.9,
        )
        return [self.classifier.classify(extracted)]

    def _deduplicate_and_classify(
        self, images: list[ExtractedImage], source_key: str
    ) -> list[ExtractedImage]:
        seen_hashes = set()
        unique_images = []
        for img in images:
            if img.file_path and os.path.exists(img.file_path):
                try:
                    import hashlib

                    with open(img.file_path, "rb") as f:
                        file_hash = hashlib.md5(f.read(), usedforsecurity=False).hexdigest()
                    img.md5_hash = file_hash
                    if file_hash in seen_hashes:
                        if os.path.exists(img.file_path):
                            os.remove(img.file_path)
                        continue
                    seen_hashes.add(file_hash)
                except Exception:
                    pass

            img = self.classifier.classify(img)
            if not img.is_decorative:
                img = self.classifier.check_template_repeat(img, source_key)
            if not img.is_decorative or img.source_shape_type in ("chart", "smartart"):
                unique_images.append(img)

        return unique_images
