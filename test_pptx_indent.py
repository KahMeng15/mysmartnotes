from pptx import Presentation

prs = Presentation("scripts/ProcessingAlgorithmTest/input/Topic03-Inheritance.pptx")
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if "Faculty extends Employee" in p.text or "public static void main" in p.text:
                    print(f"Level: {p.level}, Text: {p.text!r}")
